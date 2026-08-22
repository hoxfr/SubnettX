import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_subnettx_ppt():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # -----------------------------------------
    # SLIDE 1: Title Slide
    # -----------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SubnettX"
    subtitle.text = "Zero-Trust Dual-Vector Attendance & Threat Monitoring\n\nBuilt for Modern Infrastructure"
    
    # -----------------------------------------
    # SLIDE 2: The Problem
    # -----------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    title.text = "The Problem: Broken Trust"
    
    body = shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Traditional attendance is vulnerable to:"
    p = tf.add_paragraph()
    p.text = "Proxy Attendance ('Buddy Punching')"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "GPS Spoofing & VPN tunneling"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "QR Code screenshots sent over WhatsApp"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Conclusion: We can no longer rely on single-vector physical or digital checks."
    
    # -----------------------------------------
    # SLIDE 3: The SubnettX Solution
    # -----------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "The Solution: Dual-Vector Architecture"
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "We enforce physical presence using two simultaneous, distinct hardware vectors:"
    p = tf.add_paragraph()
    p.text = "Vector 1 (Optical): Cryptographic, rotating QR token"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Vector 2 (Acoustic): 18.5kHz Ultrasonic handshake via Web Audio API"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero-Trust Pipeline:"
    p = tf.add_paragraph()
    p.text = "If a student scans a screenshot from their dorm, the microphone will not detect the ultrasonic beacon, instantly flagging the attempt."
    p.level = 1
    
    # -----------------------------------------
    # SLIDE 4: Technology Stack
    # -----------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Technical Implementation"
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Engineered for speed, running in under 4 seconds:"
    p = tf.add_paragraph()
    p.text = "Backend: Python (Flask) for lightweight, instantaneous API routing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database: SQLite with WAL mode for ultra-fast, concurrent transactions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Frontend: Vanilla JS, Tailwind CSS, JSQR, and AudioContext for FFT frequency analysis"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Security: LEVC (Localized Endpoint Verification Core) prevents remote execution"
    p.level = 1

    # -----------------------------------------
    # SLIDE 5: Admin Threat Dashboard
    # -----------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Professor Threat Dashboard"
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Real-time command center for educators:"
    p = tf.add_paragraph()
    p.text = "Live event polling directly from the database"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Visual indicators (Green = Verified, Red = Threat Detected)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Immediate identification of students attempting to spoof presence"
    p.level = 1

    # -----------------------------------------
    # SLIDE 6: Future Scope
    # -----------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Future Roadmap"
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Scaling SubnettX for enterprise:"
    p = tf.add_paragraph()
    p.text = "Machine Learning anomaly detection on student check-in times"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Migrating SQLite to PostgreSQL for high-availability clusters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Native iOS/Android SDKs for deeper hardware sensor integration"
    p.level = 1

    # Save the presentation
    prs.save('SubnettX_Presentation.pptx')
    print("Presentation saved as SubnettX_Presentation.pptx")

if __name__ == '__main__':
    create_subnettx_ppt()
